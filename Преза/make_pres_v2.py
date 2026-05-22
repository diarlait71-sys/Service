# -*- coding: utf-8 -*-
"""
Маржа Апрель 2026 — Улучшенная презентация
Данные: апрель 2026 (строки 2026-04-01 из листа «Анализ по мес»)
"""
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

# ── ЦВЕТА ────────────────────────────────────────────────────
DARK      = RGBColor(0x0D, 0x1B, 0x2A)   # фон заголовка
MID_BLUE  = RGBColor(0x1A, 0x45, 0x7B)
ACCENT    = RGBColor(0x00, 0x78, 0xD4)
GOLD      = RGBColor(0xFF, 0xB9, 0x00)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xF3, 0xF6, 0xFA)
GRAY_LINE = RGBColor(0xD0, 0xD8, 0xE4)
GREEN     = RGBColor(0x00, 0x8A, 0x4B)
RED_BRAND = RGBColor(0xCC, 0x00, 0x00)

BRAND_CLR = {
    'KIA':       RGBColor(0x05, 0x14, 0x1F),
    'Chevrolet': RGBColor(0xB5, 0x12, 0x22),
    'JAC':       RGBColor(0x00, 0x52, 0x9B),
    'JETOUR':    RGBColor(0x00, 0x6E, 0x42),
    'Skoda':     RGBColor(0x3A, 0x75, 0xC4),
    'Soueast':   RGBColor(0x8B, 0x55, 0x00),
}

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── УТИЛИТЫ ──────────────────────────────────────────────────
def rect(slide, l, t, w, h, fill=None, line=None, lw=0.5):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h, size=14, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = 'Calibri'
    return tb

def M(n):   # миллионы с 1 знаком
    return f"{n/1_000_000:.1f}" if n else "—"

def fmt(n, d=0):
    if n is None: return "—"
    if d == 0: return f"{int(round(n)):,}".replace(",", " ")
    return f"{n:,.{d}f}".replace(",", " ")

def pct(n):
    return f"{n*100:.2f}%" if n else "—"

def header(slide, title, subtitle="", bg=DARK, stripe=GOLD):
    rect(slide, 0, 0, 13.33, 1.15, fill=bg)
    rect(slide, 0, 1.15, 13.33, 0.06, fill=stripe)
    txt(slide, title, 0.35, 0.1, 11, 0.65, size=26, bold=True)
    if subtitle:
        txt(slide, subtitle, 0.35, 0.72, 10, 0.38, size=13,
            color=RGBColor(0xAA, 0xC8, 0xEE))

def table_header(slide, cols, xs, ws, y, h, bg):
    rect(slide, xs[0], y, sum(ws)+0.02, h, fill=bg)
    for i, c in enumerate(cols):
        txt(slide, c, xs[i]+0.06, y+0.06, ws[i]-0.1, h-0.1,
            size=9, bold=True,
            align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

def table_row(slide, vals, xs, ws, y, h, bold=False, bg=None):
    if bg:
        rect(slide, xs[0], y, sum(ws)+0.02, h, fill=bg)
    for i, v in enumerate(vals):
        clr = RGBColor(0x8B, 0x50, 0x00) if bold else DARK
        txt(slide, str(v), xs[i]+0.06, y+0.05, ws[i]-0.1, h-0.08,
            size=9.5, bold=bold, color=clr,
            align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

def kpi_card(slide, x, y, w, h, value, label, val_color=GOLD, bg=DARK):
    rect(slide, x, y, w, h, fill=bg)
    rect(slide, x, y+h-0.04, w, 0.04, fill=val_color)
    txt(slide, value, x+0.08, y+0.08, w-0.16, h*0.55,
        size=22, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x+0.08, y+h*0.55, w-0.16, h*0.42,
        size=9.5, color=RGBColor(0xC0, 0xD4, 0xEE), align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 1 — ТИТУЛЬНЫЙ
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
# Диагональный акцент
rect(slide, 0, 5.6, 13.33, 0.08, fill=GOLD)
rect(slide, 0, 5.68, 13.33, 0.04, fill=ACCENT)
# Боковая полоса
rect(slide, 0.42, 1.6, 0.1, 3.5, fill=ACCENT)
# Бренд
txt(slide, "DCG GROUP", 0.75, 1.55, 10, 0.6,
    size=14, color=GOLD, bold=False)
txt(slide, "МАРЖА АПРЕЛЬ 2026", 0.75, 2.2, 11.5, 1.3,
    size=50, bold=True, color=WHITE)
txt(slide, "Анализ маржинальности по брендам и дилерским центрам",
    0.75, 3.6, 11, 0.7, size=18, color=RGBColor(0x90, 0xBB, 0xE0))
# KPI блоки на титуле
kpi_data = [
    ("2 629", "авто продано"),
    ("877 млн ₸", "маржа всего"),
    ("191 млн ₸", "доп. оборудование"),
    ("1 068 млн ₸", "к получению"),
]
for i, (v, l) in enumerate(kpi_data):
    kpi_card(slide, 0.4 + i*3.2, 4.6, 3.0, 0.9, v, l, bg=MID_BLUE)
txt(slide, "Апрель 2026 г.", 0.75, 6.0, 5, 0.4,
    size=12, color=RGBColor(0x60, 0x80, 0xA0))

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 2 — СВОДНАЯ ТАБЛИЦА (правильные данные Apr 2026)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "ИТОГИ АПРЕЛЯ 2026 — СВОДНАЯ ТАБЛИЦА",
       "Данные по всем брендам | 2 629 автомобилей")

# KPI блоки
kpi4 = [
    ("2 629", "авто", BRAND_CLR['KIA']),
    ("877 млн ₸", "маржа", BRAND_CLR['JAC']),
    ("191 млн ₸", "доп. обор.", BRAND_CLR['JETOUR']),
    ("1 068 млн ₸", "к получению", DARK),
]
for i, (v, l, bg) in enumerate(kpi4):
    kpi_card(slide, 0.25 + i*3.28, 1.38, 3.1, 0.85, v, l, bg=bg)

# Таблица брендов
brands = [
    ('KIA',       408,  632_361,  258_003_139,  66_880_100,  324_883_239, 0.0449),
    ('Chevrolet', 1787, 230_627,  412_130_000,  94_248_350,  506_378_350, 0.0358),
    ('JAC',        88,  402_827,   35_448_800,   9_576_520,   45_025_320, 0.0437),
    ('JETOUR',    343,  496_013,  170_132_400,  20_055_000,  190_187_400, 0.0473),
    ('Skoda',       3,  538_477,    1_615_430,           0,    1_615_430, 0.0218),
    ('Soueast',     2,  559_650,    1_119_300,           0,    1_119_300, 0.0352),
]
total = ('ИТОГО', 2629, 334_138, 877_329_769, 190_759_970, 1_068_089_739, None)

cols = ['Бренд', 'Авто', 'Ср. маржа ₸', 'Маржа всего', 'Доп. обор.', 'К получению', '%']
xs   = [0.25,   3.30,   4.45,           6.05,           8.05,          9.85,         12.15]
ws   = [3.05,   1.15,   1.6,            2.0,            1.8,           2.3,           1.05]
rh   = 0.52

table_header(slide, cols, xs, ws, 2.38, rh, bg=MID_BLUE)

for ri, row in enumerate(brands):
    y = 2.38 + (ri+1)*rh
    bg = LIGHT if ri % 2 == 0 else WHITE
    rect(slide, xs[0], y, sum(ws)+0.02, rh-0.04, fill=bg,
         line=GRAY_LINE, lw=0.3)
    # Цветная пометка бренда
    rect(slide, xs[0], y, 0.06, rh-0.04, fill=BRAND_CLR.get(row[0], DARK))
    vals = [row[0], fmt(row[1]), fmt(row[2]),
            f"{M(row[3])} млн", f"{M(row[4])} млн",
            f"{M(row[5])} млн", pct(row[6])]
    table_row(slide, vals, xs, ws, y, rh, bg=None)

# Итоговая строка
y = 2.38 + 7*rh
rect(slide, xs[0], y, sum(ws)+0.02, rh-0.04, fill=RGBColor(0xFF, 0xEF, 0xCC),
     line=GOLD, lw=1.0)
vals_t = [total[0], fmt(total[1]), fmt(total[2]),
          f"{M(total[3])} млн", f"{M(total[4])} млн",
          f"{M(total[5])} млн", "—"]
table_row(slide, vals_t, xs, ws, y, rh, bold=True)

# Мини-бар чарт (доля авто по брендам)
chart_data = ChartData()
chart_data.categories = [b[0] for b in brands]
chart_data.add_series('Авто', [b[1] for b in brands])

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(0.25), Inches(6.1), Inches(7.9), Inches(1.25),
    chart_data
).chart
chart.has_legend = False
chart.has_title = False
chart.plots[0].vary_by_categories = True
plot = chart.plots[0]
for i, series in enumerate(plot.series):
    for j, pt in enumerate(series.points):
        pt.format.fill.solid()
        brand = [b[0] for b in brands][j]
        pt.format.fill.fore_color.rgb = BRAND_CLR.get(brand, ACCENT)
# Y axis
chart.value_axis.has_major_gridlines = False
chart.category_axis.tick_labels.font.size = Pt(8)
txt(slide, "Количество автомобилей по брендам", 0.25, 6.0, 7.9, 0.2,
    size=8, color=MID_BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 3 — ДИНАМИКА ПО МЕСЯЦАМ (апрель 2025 → апрель 2026)
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "ДИНАМИКА ПО МЕСЯЦАМ", "Апрель 2025 — Апрель 2026")

# Данные из листа «Анализ по мес»
dyn = [
    # Метка,         Кол,  Маржа млн,  ДО млн,  К получению млн
    ('Апр\'25',      1947, 904.6,   91.5,  996.1),
    ('Май\'25',      1836, 784.7,   69.2,  853.9),
    ('Июн\'25',      2006, 840.1,   72.3,  912.4),
    ('Июл\'25',      2391, 869.8,  202.1, 1071.8),
    ('Авг\'25',      1772, 693.7,  156.8,  850.5),
    ('Сен\'25',      1835, 764.8,  200.7,  965.5),
    ('Окт\'25',      1983, 839.1,  232.2, 1071.4),
    ('Ноя\'25',      3886,1394.8,  424.0, 1818.8),
    ('Дек\'25',      2848,1022.1,  316.7, 1338.8),
    ('Янв\'26',      1809, 648.2,  202.5,  850.8),
    ('Фев\'26',      1545, 657.8,  221.1,  878.9),
    ('Мар\'26',      1505, 701.6,  219.6,  921.2),
    ('Апр\'26',      2629, 877.3,  190.8, 1068.1),
]

# Таблица
cols_d = ['Месяц', 'Авто', 'Маржа, млн ₸', 'ДО, млн ₸', 'К получению, млн ₸']
xs_d = [0.2, 2.2, 3.6, 5.5, 7.1]
ws_d = [2.0, 1.4, 1.9, 1.6, 2.2]
rh_d = 0.40

table_header(slide, cols_d, xs_d, ws_d, 1.32, rh_d+0.06, bg=MID_BLUE)

for ri, row in enumerate(dyn):
    y = 1.32 + (ri+1)*(rh_d+0.06)
    is_apr26 = row[0] == "Апр'26"
    is_apr25 = row[0] == "Апр'25"
    bg_clr = RGBColor(0xFF, 0xEF, 0xCC) if is_apr26 else (
             RGBColor(0xE8, 0xF0, 0xFE) if is_apr25 else
             (LIGHT if ri % 2 == 0 else WHITE))
    rect(slide, xs_d[0], y, sum(ws_d)+0.02, rh_d, fill=bg_clr,
         line=GRAY_LINE, lw=0.3)
    vals = [row[0], fmt(row[1]), f"{row[2]:.1f}", f"{row[3]:.1f}", f"{row[4]:.1f}"]
    table_row(slide, vals, xs_d, ws_d, y, rh_d, bold=is_apr26)

# Легенда
txt(slide, "★ Апрель 2026", xs_d[0]+0.1, 1.32 + 14*(rh_d+0.06)+0.05, 3, 0.3,
    size=8, color=RGBColor(0x8B, 0x50, 0x00), bold=True)

# Колоночная диаграмма кол-во авто
chart_d = ChartData()
chart_d.categories = [r[0] for r in dyn]
chart_d.add_series('Авто', [r[1] for r in dyn])

ch = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(9.3), Inches(1.32), Inches(3.85), Inches(4.5),
    chart_d
).chart
ch.has_legend = False
ch.has_title = True
ch.chart_title.text_frame.text = "Кол-во авто"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
ch.value_axis.has_major_gridlines = True
ch.category_axis.tick_labels.font.size = Pt(7)
ch.value_axis.tick_labels.font.size = Pt(8)

# Цвет последнего столбца (апр'26)
plot_d = ch.plots[0]
plot_d.vary_by_categories = True
for series in plot_d.series:
    for j, pt in enumerate(series.points):
        pt.format.fill.solid()
        if j == len(dyn) - 1:  # апрель 2026 — золотой
            pt.format.fill.fore_color.rgb = GOLD
        elif j == 0:  # апрель 2025 — синий
            pt.format.fill.fore_color.rgb = ACCENT
        else:
            pt.format.fill.fore_color.rgb = MID_BLUE

# Маржа chart
chart_m = ChartData()
chart_m.categories = [r[0] for r in dyn]
chart_m.add_series('Маржа (млн ₸)', [r[2] for r in dyn])

ch2 = slide.shapes.add_chart(
    XL_CHART_TYPE.LINE,
    Inches(9.3), Inches(5.9), Inches(3.85), Inches(1.45),
    chart_m
).chart
ch2.has_legend = False
ch2.has_title = True
ch2.chart_title.text_frame.text = "Маржа, млн ₸"
ch2.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
ch2.category_axis.tick_labels.font.size = Pt(7)
ch2.value_axis.tick_labels.font.size = Pt(7)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 4 — KIA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "KIA — ДИЛЕРЫ И МОДЕЛИ",
       "Апрель 2026 | 408 авто | Маржа: 258 млн ₸ | 4.49%",
       bg=BRAND_CLR['KIA'], stripe=GOLD)

kia_dealers = [
    ('Актау',        18, 679_219, 12_225_939,  2_942_000, 0.0483),
    ('Алматы Qulja', 57, 614_007, 34_998_400,  8_300_000, 0.0419),
    ('Шымкент',     125, 681_154, 85_144_200, 29_148_100, 0.0485),
    ('Алматы ALMATY',59, 569_750, 33_615_250,  6_050_000, 0.0413),
    ('Тараз',        23, 592_709, 13_632_300,  7_650_000, 0.0426),
    ('Астана',       65, 625_837, 40_679_400,  8_690_000, 0.0436),
    ('УКО',          34, 536_531, 18_242_050,  4_100_000, 0.0403),
    ('Павлодар',     27, 720_948, 19_465_600,          0, 0.0516),
    ('ИТОГО',       408, 632_361,258_003_139, 66_880_100, 0.0449),
]

kia_models = [
    ('K5',          91, 675_655, 0.0502, 22.3),
    ('Sportage',    93, 788_217, 0.0492, 22.8),
    ('Sorento 2025',61, 743_283, 0.0395, 14.9),
    ('Cerato',      48, 382_550, 0.0367, 11.8),
    ('Soluto',      38, 298_947, 0.0405,  9.3),
    ('Seltos',      26, 558_340, 0.0448,  6.4),
    ('Sorento 2026',11,1_123_314,0.0502,  2.7),
    ('Ceed SW',     10, 512_550, 0.0450,  2.5),
    ('Carnival',     3,1_234_417,0.0508,  0.7),
    ('Ceed',         6, 437_642, 0.0443,  1.5),
    ('K8',           1,1_721_850,0.0650,  0.2),
    ('Soul',         4, 540_675, 0.0450,  1.0),
    ('Sportage C',  16, 371_200, 0.0290,  3.9),
]

# Левая часть — дилеры
cols_k = ['Дилер', 'Авто', 'Ср.маржа', 'М.пост.скид', 'ДОП', '%']
xs_k   = [0.2, 3.4, 4.75, 6.3, 8.15, 9.5]
ws_k   = [3.2, 1.35, 1.55, 1.85, 1.35, 0.85]
rh_k   = 0.49

table_header(slide, cols_k, xs_k, ws_k, 1.32, rh_k, bg=BRAND_CLR['KIA'])

for ri, row in enumerate(kia_dealers):
    y = 1.32 + (ri+1)*rh_k
    is_tot = row[0] == 'ИТОГО'
    bg_c = RGBColor(0xFF, 0xEF, 0xCC) if is_tot else (LIGHT if ri%2==0 else WHITE)
    rect(slide, xs_k[0], y, sum(ws_k)+0.02, rh_k-0.04, fill=bg_c,
         line=GRAY_LINE, lw=0.3)
    if not is_tot:
        rect(slide, xs_k[0], y, 0.06, rh_k-0.04, fill=BRAND_CLR['KIA'])
    vals = [row[0], str(row[1]), fmt(row[2]),
            f"{M(row[3])} млн", f"{M(row[4])} млн", pct(row[5])]
    table_row(slide, vals, xs_k, ws_k, y, rh_k, bold=is_tot)

# Правая часть — модели (чарт)
model_names = [m[0] for m in kia_models]
model_vals  = [m[1] for m in kia_models]

cd_k = ChartData()
cd_k.categories = model_names
cd_k.add_series('Авто', model_vals)

ch_k = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(10.5), Inches(1.32), Inches(2.65), Inches(5.1),
    cd_k
).chart
ch_k.has_legend = False
ch_k.has_title = True
ch_k.chart_title.text_frame.text = "Продажи по моделям"
ch_k.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
ch_k.category_axis.tick_labels.font.size = Pt(7)
ch_k.value_axis.tick_labels.font.size = Pt(7)
ch_k.value_axis.has_major_gridlines = False
for s in ch_k.plots[0].series:
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = BRAND_CLR['KIA']

# Итог KIA внизу
rect(slide, 0.2, 6.55, 10.15, 0.75, fill=BRAND_CLR['KIA'])
kia_footer = [
    ("Маржа инвест", "28.7 млн ₸", 0.3),
    ("Скидки (15 авто)", "12.3 млн ₸", 3.0),
    ("Лидер по авто", "Шымкент (125)", 5.7),
    ("Топ % маржи", "Павлодар 5.16%", 8.0),
]
for lbl, val, x in kia_footer:
    txt(slide, lbl,  x, 6.58, 2.5, 0.28, size=8.5, color=RGBColor(0xAA, 0xCC, 0xFF))
    txt(slide, val,  x, 6.86, 2.5, 0.35, size=12, bold=True, color=GOLD)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 5 — CHEVROLET
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "CHEVROLET — ДИЛЕРЫ И МОДЕЛИ",
       "Апрель 2026 | 1 787 авто | Маржа: 412 млн ₸ | 3.58%",
       bg=BRAND_CLR['Chevrolet'], stripe=GOLD)

chev_dealers = [
    ('Актау',         91, 227_473,  20_700_000,  8_430_000, 0.0373),
    ('Атырау',        72, 237_500,  17_100_000,  2_350_000, 0.0387),
    ('Алматы 1159',    4, 300_000,   1_200_000,    900_000, 0.0443),
    ('Шымкент 119',  375, 219_573,  82_340_000, 15_165_000, 0.0345),
    ('Алматы Zhets.',  72, 242_361,  17_450_000,          0, 0.0333),
    ('Тараз',        118, 219_915,  25_950_000,  6_850_000, 0.0357),
    ('Астана',       108, 268_519,  29_000_000,  3_897_000, 0.0360),
    ('УКО',          130, 226_923,  29_500_000, 12_140_000, 0.0369),
    ('Павлодар',      94, 245_745,  23_100_000,  4_770_000, 0.0367),
    ('Уральск',      124, 234_274,  29_050_000,  3_201_850, 0.0375),
    ('Кызылорда',    130, 220_769,  28_700_000,  5_890_000, 0.0359),
    ('Шымкент 935',  469, 230_362, 108_040_000, 30_654_500, 0.0355),
    ('ИТОГО',       1787, 230_627, 412_130_000, 94_248_350, 0.0358),
]

chev_models = [
    ('Cobalt',    1748, 223_587, 0.0364, 97.8),
    ('Onix',        23, 330_435, 0.0432,  1.3),
    ('Tahoe MY25',  11,1_118_182,0.0220,  0.6),
    ('Tracker',      2, 400_000, 0.0381,  0.1),
    ('Labo',         3, 200_000, 0.0358,  0.2),
]

cols_c = ['Дилер', 'Авто', 'Ср.маржа', 'Маржа', 'ДОП', '%']
xs_c = [0.2, 3.4, 4.75, 6.3, 8.15, 9.5]
ws_c = [3.2, 1.35, 1.55, 1.85, 1.35, 0.85]
rh_c = 0.43

table_header(slide, cols_c, xs_c, ws_c, 1.32, rh_c, bg=BRAND_CLR['Chevrolet'])

for ri, row in enumerate(chev_dealers):
    y = 1.32 + (ri+1)*rh_c
    is_tot = row[0]=='ИТОГО'
    bg_c = RGBColor(0xFF, 0xEF, 0xCC) if is_tot else (LIGHT if ri%2==0 else WHITE)
    rect(slide, xs_c[0], y, sum(ws_c)+0.02, rh_c-0.03, fill=bg_c,
         line=GRAY_LINE, lw=0.3)
    if not is_tot:
        rect(slide, xs_c[0], y, 0.06, rh_c-0.03, fill=BRAND_CLR['Chevrolet'])
    vals = [row[0], str(row[1]), fmt(row[2]),
            f"{M(row[3])} млн", f"{M(row[4])} млн", pct(row[5])]
    table_row(slide, vals, xs_c, ws_c, y, rh_c, bold=is_tot)

# Модели — чарт
cd_cv = ChartData()
cd_cv.categories = [m[0] for m in chev_models]
cd_cv.add_series('Авто', [m[1] for m in chev_models])

ch_cv = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE,
    Inches(10.5), Inches(1.32), Inches(2.65), Inches(3.5),
    cd_cv
).chart
ch_cv.has_legend = True
ch_cv.has_title = True
ch_cv.chart_title.text_frame.text = "Доля моделей"
ch_cv.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
if ch_cv.legend:
    ch_cv.legend.font.size = Pt(8)

# Таблица моделей
col_cm = ['Модель', 'Авто', 'Ср.маржа', '%']
xs_cm  = [10.5, 11.85, 12.45, 12.95]
ws_cm  = [1.35, 0.6, 0.5, 0.4]
table_header(slide, col_cm, xs_cm, ws_cm, 5.0, 0.42, bg=BRAND_CLR['Chevrolet'])
for ri, m in enumerate(chev_models):
    y = 5.0 + (ri+1)*0.42
    bg_c = LIGHT if ri%2==0 else WHITE
    rect(slide, xs_cm[0], y, sum(ws_cm), 0.38, fill=bg_c, line=GRAY_LINE, lw=0.3)
    table_row(slide, [m[0], str(m[1]), fmt(m[2]), f"{m[4]:.1f}%"],
              xs_cm, ws_cm, y, 0.42)

# Footer
rect(slide, 0.2, 6.72, 10.15, 0.6, fill=BRAND_CLR['Chevrolet'])
c_footer = [
    ("Лидер", "Шымкент 935 (469 авто)", 0.3),
    ("Скидки (48 авто)", "4.2 млн ₸", 4.5),
    ("Топ % маржи", "Алматы 4.43%", 7.5),
]
for lbl, val, x in c_footer:
    txt(slide, lbl, x, 6.74, 3, 0.25, size=8.5, color=RGBColor(0xFF, 0xBB, 0xBB))
    txt(slide, val, x, 6.99, 3, 0.28, size=11, bold=True, color=GOLD)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 6 — JAC
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "JAC — ДИЛЕРЫ И МОДЕЛИ",
       "Апрель 2026 | 88 авто | Маржа: 35.4 млн ₸ | 4.37%",
       bg=BRAND_CLR['JAC'], stripe=GOLD)

jac_dealers = [
    ('Шымкент Baidybek', 14, 466_843,  6_535_800,  2_280_000, 0.0477),
    ('Тараз',             2, 475_600,    951_200,           0, 0.0400),
    ('Актау',             3, 520_200,  1_560_600,           0, 0.0476),
    ('Атырау',            5, 446_440,  2_232_200,     30_000, 0.0431),
    ('Шымкент ONTUSTIK', 11, 306_000,  3_366_000,    500_000, 0.0364),
    ('УКО',              12, 423_733,  5_084_800,  2_350_000, 0.0458),
    ('Шымкент Темир.',    7, 303_600,  2_125_200,  1_250_000, 0.0400),
    ('Астана',           12, 383_200,  4_598_400,  2_058_720, 0.0417),
    ('Павлодар',          9, 465_778,  4_192_000,    535_000, 0.0480),
    ('Кызылорда',         6, 334_933,  2_009_600,    200_000, 0.0400),
    ('Уральск',           7, 399_000,  2_793_000,    372_800, 0.0451),
    ('ИТОГО',            88, 402_827, 35_448_800,  9_576_520, 0.0437),
]

jac_models = [
    ('S3 Pro', 40, 274_200, 0.0400, 45.5),
    ('J7',     18, 311_489, 0.0398, 20.5),
    ('T9',      9, 971_756, 0.0558, 10.2),
    ('JS8',     7, 475_600, 0.0400,  8.0),
    ('JS4',     5, 350_000, 0.0400,  5.7),
    ('T6',      5, 513_680, 0.0422,  5.7),
    ('T8 Pro',  2, 828_700, 0.0526,  2.3),
    ('J7 PLUS', 2, 411_600, 0.0400,  2.3),
]

cols_j = ['Дилер', 'Авто', 'Ср.маржа', 'Маржа', 'ДОП', '%']
xs_j = [0.2, 3.85, 5.2, 6.65, 8.3, 9.6]
ws_j = [3.65, 1.35, 1.45, 1.65, 1.3, 0.8]
rh_j = 0.475

table_header(slide, cols_j, xs_j, ws_j, 1.32, rh_j, bg=BRAND_CLR['JAC'])
for ri, row in enumerate(jac_dealers):
    y = 1.32 + (ri+1)*rh_j
    is_tot = row[0]=='ИТОГО'
    bg_c = RGBColor(0xFF, 0xEF, 0xCC) if is_tot else (LIGHT if ri%2==0 else WHITE)
    rect(slide, xs_j[0], y, sum(ws_j)+0.02, rh_j-0.04, fill=bg_c,
         line=GRAY_LINE, lw=0.3)
    if not is_tot:
        rect(slide, xs_j[0], y, 0.06, rh_j-0.04, fill=BRAND_CLR['JAC'])
    vals = [row[0], str(row[1]), fmt(row[2]),
            f"{M(row[3])} млн", f"{M(row[4])} млн", pct(row[5])]
    table_row(slide, vals, xs_j, ws_j, y, rh_j, bold=is_tot)

# Модели — чарт
cd_j = ChartData()
cd_j.categories = [m[0] for m in jac_models]
cd_j.add_series('Авто', [m[1] for m in jac_models])

ch_j = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(10.5), Inches(1.32), Inches(2.65), Inches(4.0),
    cd_j
).chart
ch_j.has_legend = False
ch_j.has_title = True
ch_j.chart_title.text_frame.text = "Авто по моделям"
ch_j.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
ch_j.category_axis.tick_labels.font.size = Pt(7)
ch_j.value_axis.tick_labels.font.size = Pt(7)
ch_j.value_axis.has_major_gridlines = False
for s in ch_j.plots[0].series:
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = BRAND_CLR['JAC']

# ср маржа чарт
cd_j2 = ChartData()
cd_j2.categories = [m[0] for m in jac_models]
cd_j2.add_series('Ср. маржа', [m[2] for m in jac_models])

ch_j2 = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(10.5), Inches(5.5), Inches(2.65), Inches(1.8),
    cd_j2
).chart
ch_j2.has_legend = False
ch_j2.has_title = True
ch_j2.chart_title.text_frame.text = "Ср. маржа ₸"
ch_j2.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
ch_j2.category_axis.tick_labels.font.size = Pt(7)
ch_j2.value_axis.tick_labels.font.size = Pt(7)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 7 — JETOUR
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "JETOUR — ДИЛЕРЫ И МОДЕЛИ",
       "Апрель 2026 | 343 авто | Маржа: 170 млн ₸ | 4.73%",
       bg=BRAND_CLR['JETOUR'], stripe=GOLD)

jet_dealers = [
    ('Алматы Alatau',  26, 495_850, 12_892_100,  1_250_000, 0.0440),
    ('Актау',          32, 634_025, 20_288_800,  2_500_000, 0.0519),
    ('Атырау',         34, 369_285, 12_555_700,  5_395_000, 0.0411),
    ('Алматы Zhets.',  31, 552_371, 17_123_500,           0, 0.0484),
    ('УКО',            50, 513_714, 25_685_700,  2_750_000, 0.0505),
    ('Шымкент Ont.',   33, 499_097, 16_470_200,  2_250_000, 0.0452),
    ('Павлодар',       29, 443_338, 12_856_800,    250_000, 0.0462),
    ('Кызылорда',      22, 428_186,  9_420_100,  1_400_000, 0.0449),
    ('Тараз',           8, 472_025,  3_776_200,  1_250_000, 0.0475),
    ('Шымкент Baid.',  29, 574_697, 16_666_200,  2_010_000, 0.0488),
    ('Уральск',        49, 457_084, 22_397_100,  1_000_000, 0.0477),
    ('ИТОГО',         343, 496_013,170_132_400, 20_055_000, 0.0473),
]

jet_models = [
    ('X70FL',   125, 417_673, 0.0479, 36.4),
    ('X50',      84, 299_600, 0.0392, 24.5),
    ('X70Plus',  49, 548_947, 0.0443, 14.3),
    ('T1',       30, 927_733, 0.0589,  8.7),
    ('T2',       29, 896_103, 0.0549,  8.5),
    ('X90Plus',  16, 519_600, 0.0395,  4.7),
    ('Dashing',   9, 364_044, 0.0360,  2.6),
    ('X70',       1, 449_500, 0.0500,  0.3),
]

cols_jt = ['Дилер', 'Авто', 'Ср.маржа', 'Маржа', 'ДОП', '%']
xs_jt = [0.2, 3.85, 5.2, 6.65, 8.3, 9.6]
ws_jt = [3.65, 1.35, 1.45, 1.65, 1.3, 0.8]
rh_jt = 0.475

table_header(slide, cols_jt, xs_jt, ws_jt, 1.32, rh_jt, bg=BRAND_CLR['JETOUR'])
for ri, row in enumerate(jet_dealers):
    y = 1.32 + (ri+1)*rh_jt
    is_tot = row[0]=='ИТОГО'
    bg_c = RGBColor(0xFF, 0xEF, 0xCC) if is_tot else (LIGHT if ri%2==0 else WHITE)
    rect(slide, xs_jt[0], y, sum(ws_jt)+0.02, rh_jt-0.04, fill=bg_c,
         line=GRAY_LINE, lw=0.3)
    if not is_tot:
        rect(slide, xs_jt[0], y, 0.06, rh_jt-0.04, fill=BRAND_CLR['JETOUR'])
    vals = [row[0], str(row[1]), fmt(row[2]),
            f"{M(row[3])} млн", f"{M(row[4])} млн", pct(row[5])]
    table_row(slide, vals, xs_jt, ws_jt, y, rh_jt, bold=is_tot)

# Модели — чарт
cd_jt = ChartData()
cd_jt.categories = [m[0] for m in jet_models]
cd_jt.add_series('Авто', [m[1] for m in jet_models])

ch_jt = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED,
    Inches(10.5), Inches(1.32), Inches(2.65), Inches(4.0),
    cd_jt
).chart
ch_jt.has_legend = False
ch_jt.has_title = True
ch_jt.chart_title.text_frame.text = "Авто по моделям"
ch_jt.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(9)
ch_jt.category_axis.tick_labels.font.size = Pt(7)
ch_jt.value_axis.tick_labels.font.size = Pt(7)
ch_jt.value_axis.has_major_gridlines = False
for s in ch_jt.plots[0].series:
    s.format.fill.solid()
    s.format.fill.fore_color.rgb = BRAND_CLR['JETOUR']

# Топ модели
txt(slide, "Топ по авто: X70FL (125), X50 (84)", 10.55, 5.45, 2.6, 0.3,
    size=8, color=DARK, bold=True)
txt(slide, "Топ по марже: T1 (5.89%), T2 (5.49%)", 10.55, 5.75, 2.6, 0.3,
    size=8, color=DARK, bold=True)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 8 — SKODA + SOUEAST + СРАВНЕНИЕ БРЕНДОВ
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT)
header(slide, "SKODA | SOUEAST — НОВЫЕ БРЕНДЫ + СРАВНЕНИЕ",
       "Апрель 2026")

# Skoda
rect(slide, 0.2, 1.32, 4.4, 0.42, fill=BRAND_CLR['Skoda'])
txt(slide, "SKODA — 3 авто", 0.3, 1.36, 4.2, 0.34, size=12, bold=True)

skoda_rows = [
    ('УКО',    1, 465_430, 0.0189),
    ('Атырау', 2, 575_000, 0.0233),
    ('ИТОГО',  3, 538_477, 0.0218),
]
cols_sk = ['Дилер', 'Авто', 'Ср. маржа', '%']
xs_sk = [0.2, 1.8, 2.7, 4.0]
ws_sk = [1.6, 0.9, 1.3, 0.7]
table_header(slide, cols_sk, xs_sk, ws_sk, 1.74, 0.42, bg=BRAND_CLR['Skoda'])
for ri, r in enumerate(skoda_rows):
    y = 1.74 + (ri+1)*0.42
    bg_c = RGBColor(0xFF, 0xEF, 0xCC) if r[0]=='ИТОГО' else (LIGHT if ri%2==0 else WHITE)
    rect(slide, xs_sk[0], y, sum(ws_sk), 0.38, fill=bg_c, line=GRAY_LINE, lw=0.3)
    table_row(slide, [r[0], str(r[1]), fmt(r[2]), pct(r[3])],
              xs_sk, ws_sk, y, 0.42, bold=r[0]=='ИТОГО')

txt(slide, "Модель: Kodiaq NG | % маржи: 2.18%", 0.3, 3.12, 4.2, 0.4,
    size=9, color=DARK)

# Soueast
rect(slide, 0.2, 3.7, 4.4, 0.42, fill=BRAND_CLR['Soueast'])
txt(slide, "SOUEAST — 2 авто", 0.3, 3.74, 4.2, 0.34, size=12, bold=True)
rect(slide, 0.2, 4.12, 4.4, 0.5, fill=LIGHT, line=GRAY_LINE, lw=0.5)
txt(slide, "Ср. маржа: 559 650 ₸  |  К получению: 1.1 млн ₸  |  % маржи: 3.52%",
    0.3, 4.18, 4.2, 0.38, size=9.5, color=DARK)

# Сравнение брендов — столбчатый
all_brands_cmp = [
    ('KIA',       408,  632_361, 0.0449),
    ('Chevrolet', 1787, 230_627, 0.0358),
    ('JAC',        88,  402_827, 0.0437),
    ('JETOUR',    343,  496_013, 0.0473),
    ('Skoda',       3,  538_477, 0.0218),
    ('Soueast',     2,  559_650, 0.0352),
]

# Ср маржа по брендам
cd_cmp = ChartData()
cd_cmp.categories = [b[0] for b in all_brands_cmp]
cd_cmp.add_series('Ср. маржа ₸', [b[2] for b in all_brands_cmp])

ch_cmp = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(4.9), Inches(1.32), Inches(4.0), Inches(3.2),
    cd_cmp
).chart
ch_cmp.has_legend = False
ch_cmp.has_title = True
ch_cmp.chart_title.text_frame.text = "Средняя маржа по брендам, ₸"
ch_cmp.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
ch_cmp.category_axis.tick_labels.font.size = Pt(9)
ch_cmp.value_axis.tick_labels.font.size = Pt(8)
ch_cmp.plots[0].vary_by_categories = True
for s in ch_cmp.plots[0].series:
    for j, pt in enumerate(s.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = BRAND_CLR.get(all_brands_cmp[j][0], ACCENT)

# % маржи по брендам
cd_pct = ChartData()
cd_pct.categories = [b[0] for b in all_brands_cmp]
cd_pct.add_series('% маржи', [b[3]*100 for b in all_brands_cmp])

ch_pct = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED,
    Inches(9.1), Inches(1.32), Inches(4.0), Inches(3.2),
    cd_pct
).chart
ch_pct.has_legend = False
ch_pct.has_title = True
ch_pct.chart_title.text_frame.text = "% маржи по брендам"
ch_pct.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(10)
ch_pct.category_axis.tick_labels.font.size = Pt(9)
ch_pct.value_axis.tick_labels.font.size = Pt(8)
ch_pct.plots[0].vary_by_categories = True
for s in ch_pct.plots[0].series:
    for j, pt in enumerate(s.points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = BRAND_CLR.get(all_brands_cmp[j][0], ACCENT)

# Сводная таблица сравнения
cols_b = ['Бренд', 'Авто', 'Ср.маржа', '%']
xs_b = [4.9, 7.2, 8.5, 9.5]
ws_b = [2.3, 1.3, 1.0, 0.7]
table_header(slide, cols_b, xs_b, ws_b, 4.7, 0.42, bg=MID_BLUE)
for ri, b in enumerate(all_brands_cmp):
    y = 4.7 + (ri+1)*0.42
    bg_c = LIGHT if ri%2==0 else WHITE
    rect(slide, xs_b[0], y, sum(ws_b), 0.38, fill=bg_c, line=GRAY_LINE, lw=0.3)
    rect(slide, xs_b[0], y, 0.06, 0.38, fill=BRAND_CLR.get(b[0], DARK))
    table_row(slide, [b[0], fmt(b[1]), fmt(b[2]), pct(b[3])],
              xs_b, ws_b, y, 0.42)

# ═══════════════════════════════════════════════════════════════
# СЛАЙД 9 — ИТОГИ И ВЫВОДЫ
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
rect(slide, 0, 0, 13.33, 7.5, fill=DARK)
rect(slide, 0, 0, 13.33, 0.08, fill=GOLD)
rect(slide, 0, 7.42, 13.33, 0.08, fill=GOLD)
rect(slide, 0.4, 1.2, 0.08, 4.8, fill=ACCENT)

txt(slide, "КЛЮЧЕВЫЕ ВЫВОДЫ", 0.7, 1.15, 11, 0.5,
    size=13, color=GOLD)
txt(slide, "АПРЕЛЬ 2026", 0.7, 1.65, 11, 0.9,
    size=44, bold=True, color=WHITE)

# Разделитель
rect(slide, 0.7, 2.6, 11.5, 0.04, fill=ACCENT)

conclusions = [
    ("🚗  Итого:", "2 629 авто продано | Маржа к получению: 1 068 млн ₸"),
    ("🏆  Лидер по объёму:", "Chevrolet — 1 787 авто (68% всех продаж)"),
    ("💎  Лидер по ср. марже:", "KIA — 632 360 ₸/авто (4.49%) | 324.9 млн к получению"),
    ("📈  Лидер по % маржи:", "JETOUR — 4.73% | Топ модели: T1 (5.89%), T2 (5.49%)"),
    ("📊  JAC:", "88 авто, 4.37% — стабильный результат"),
    ("🆕  Новые бренды:", "Skoda (3 авто, 2.18%) + Soueast (2 авто, 3.52%)"),
    ("📅  Динамика:", "Апрель 2026 > апреля 2025 на 35% по кол-ву авто (2629 vs 1947)"),
]

for i, (bold_part, normal_part) in enumerate(conclusions):
    y = 2.78 + i * 0.57
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(12.3), Inches(0.52))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = bold_part + "  "
    r1.font.size = Pt(14); r1.font.bold = True
    r1.font.color.rgb = GOLD; r1.font.name = 'Calibri'
    r2 = p.add_run(); r2.text = normal_part
    r2.font.size = Pt(14); r2.font.bold = False
    r2.font.color.rgb = WHITE; r2.font.name = 'Calibri'

# ── СОХРАНИТЬ ────────────────────────────────────────────────
out = r"c:\Users\D.Muldabaev\Desktop\Приложения для сервиса\Преза\Маржа_Апрель_2026_v2.pptx"
prs.save(out)
print(f"Готово: {out}")
