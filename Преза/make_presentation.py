# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Цвета бренда ──────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x2B, 0x55)   # тёмно-синий фон
ACCENT_BLUE = RGBColor(0x00, 0x70, 0xC0)   # акцентный синий
LIGHT_GRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # светлый фон таблицы
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GOLD        = RGBColor(0xFF, 0xC0, 0x00)
GREEN       = RGBColor(0x00, 0x97, 0x60)
ORANGE      = RGBColor(0xE8, 0x6A, 0x10)

# ── Цвета брендов ─────────────────────────────────────────────
BRAND_COLORS = {
    'KIA':        RGBColor(0x05, 0x14, 0x1E),
    'Chevrolet':  RGBColor(0xD4, 0x1B, 0x2C),
    'JAC':        RGBColor(0x00, 0x5B, 0xAA),
    'JETOUR':     RGBColor(0x00, 0x7A, 0x4E),
    'Skoda':      RGBColor(0x4A, 0x90, 0xD9),
}

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank

# ═══════════════════════════════════════════════════════════════
# Вспомогательные функции
# ═══════════════════════════════════════════════════════════════

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_w=None):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def fmt_num(n, decimals=0):
    """Форматирование числа с разрядами"""
    if n is None:
        return '—'
    if decimals == 0:
        return f"{int(round(n)):,}".replace(',', ' ')
    return f"{n:,.{decimals}f}".replace(',', ' ')

def fmt_mln(n):
    """Число в миллионах"""
    if n is None:
        return '—'
    return f"{n/1_000_000:.1f} млн"

def fmt_pct(n):
    if n is None:
        return '—'
    return f"{n*100:.1f}%"

# ═══════════════════════════════════════════════════════════════
# Слайд 1 — Титульный
# ═══════════════════════════════════════════════════════════════
def slide_title():
    slide = prs.slides.add_slide(BLANK)
    # Фон
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_BLUE)
    # Золотая полоса
    add_rect(slide, 0, 5.8, 13.33, 0.12, fill_color=GOLD)
    # Акцентная вертикальная линия
    add_rect(slide, 0.5, 1.8, 0.07, 3.2, fill_color=ACCENT_BLUE)

    add_text(slide, 'DCG GROUP', 0.8, 1.7, 12, 1,
             font_size=16, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(slide, 'МАРЖА АПРЕЛЬ 2026', 0.8, 2.3, 11.5, 1.5,
             font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, 'Анализ маржинальности по брендам и дилерским центрам', 0.8, 3.9, 11, 0.8,
             font_size=20, bold=False, color=RGBColor(0xB0, 0xC8, 0xE8), align=PP_ALIGN.LEFT)
    add_text(slide, 'Апрель 2026 г.', 0.8, 6.1, 4, 0.5,
             font_size=14, bold=False, color=RGBColor(0xAA, 0xAA, 0xAA), align=PP_ALIGN.LEFT)

slide_title()

# ═══════════════════════════════════════════════════════════════
# Слайд 2 — Итоги апреля сводная таблица
# ═══════════════════════════════════════════════════════════════
def slide_summary():
    slide = prs.slides.add_slide(BLANK)
    # Шапка
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=DARK_BLUE)
    add_text(slide, 'ИТОГИ АПРЕЛЯ 2026 — СВОДНАЯ ТАБЛИЦА', 0.3, 0.15, 12, 0.8,
             font_size=24, bold=True, color=WHITE)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    # Данные
    headers = ['Бренд', 'Кол-во авто', 'Ср. маржа', 'Маржа всего', 'ДО', 'К получению']
    rows = [
        ('KIA',       463,  605417,  280_308_275,  84_053_500,  364_361_775),
        ('Chevrolet', 858,  264234,  226_712_827,     950_000,  227_662_827),
        ('JAC',       190,  578659,  109_945_150,   2_550_000,  112_495_150),
        ('JETOUR',    436,  659702,  287_630_100,   3_992_200,  291_622_300),
        ('ИТОГО',    1947,  464610,  904_596_352,  91_545_700,  996_142_052),
    ]

    col_x = [0.2, 2.2, 4.2, 6.3, 9.0, 10.9]
    col_w = [2.0, 2.0, 2.1, 2.7, 1.9, 2.3]
    row_h = 0.72

    # Заголовок таблицы
    add_rect(slide, 0.2, 1.3, 12.9, row_h, fill_color=ACCENT_BLUE)
    for i, h in enumerate(headers):
        add_text(slide, h, col_x[i], 1.35, col_w[i], row_h - 0.1,
                 font_size=13, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    brand_colors_row = {
        'KIA': RGBColor(0xE8, 0xF0, 0xFE),
        'Chevrolet': RGBColor(0xFD, 0xED, 0xED),
        'JAC': RGBColor(0xE8, 0xF4, 0xFF),
        'JETOUR': RGBColor(0xE8, 0xF8, 0xF0),
        'ИТОГО': RGBColor(0xFF, 0xF0, 0xCC),
    }

    for ri, row in enumerate(rows):
        y = 1.3 + (ri + 1) * row_h
        bg = brand_colors_row.get(row[0], LIGHT_GRAY)
        add_rect(slide, 0.2, y, 12.9, row_h - 0.04, fill_color=bg)
        values = [
            row[0],
            fmt_num(row[1]),
            fmt_num(row[2]),
            fmt_mln(row[3]),
            fmt_mln(row[4]),
            fmt_mln(row[5]),
        ]
        bold_row = row[0] == 'ИТОГО'
        for ci, val in enumerate(values):
            add_text(slide, val, col_x[ci], y + 0.08, col_w[ci], row_h - 0.15,
                     font_size=13, bold=bold_row,
                     color=DARK_BLUE if not bold_row else RGBColor(0x8B, 0x60, 0x00),
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # KPI блоки внизу
    kpis = [
        ('1 947', 'авто продано'),
        ('905 млн ₸', 'маржа всего'),
        ('91.5 млн ₸', 'доп. оборудование'),
        ('996 млн ₸', 'к получению'),
    ]
    kpi_x = [0.3, 3.5, 6.7, 9.9]
    for i, (val, lbl) in enumerate(kpis):
        add_rect(slide, kpi_x[i], 6.2, 3.0, 1.0, fill_color=DARK_BLUE)
        add_text(slide, val, kpi_x[i], 6.2, 3.0, 0.55,
                 font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_text(slide, lbl, kpi_x[i], 6.7, 3.0, 0.45,
                 font_size=11, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

slide_summary()

# ═══════════════════════════════════════════════════════════════
# Слайд 3 — KIA анализ по дилерам
# ═══════════════════════════════════════════════════════════════
def slide_kia_dealers():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BRAND_COLORS['KIA'])
    add_text(slide, 'KIA — АНАЛИЗ ПО ДИЛЕРСКИМ ЦЕНТРАМ', 0.3, 0.15, 12, 0.8,
             font_size=22, bold=True, color=WHITE)
    add_text(slide, 'Апрель 2026', 11.0, 0.18, 2.0, 0.5,
             font_size=13, bold=False, color=GOLD, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    headers = ['Дилерский центр', 'Кол', 'РРЦ', 'Маржа', 'Скидка', 'М. после скидок', 'ДОП', 'Ср. маржа', '%']
    dealers = [
        ('Актау',             18,  14_054_433,  12_585_739,    359_800,  12_225_939,  2_942_000,   679_219, 0.0483),
        ('Алматы Qulja',      57,  14_670_702,  41_978_400,  6_980_000,  34_998_400,  8_300_000,   614_007, 0.0419),
        ('Шымкент ONTUSTIK', 125,  14_035_600,  85_360_200,    216_000,  85_144_200, 29_148_100,   681_154, 0.0485),
        ('Алматы ALMATY',     59,  13_784_915,  35_633_650,  2_018_400,  33_615_250,  6_050_000,   569_750, 0.0413),
        ('Тараз',             23,  13_916_087,  14_232_300,    600_000,  13_632_300,  7_650_000,   592_709, 0.0426),
        ('Астана',            65,  14_362_308,  40_969_400,    290_000,  40_679_400,  8_690_000,   625_837, 0.0436),
        ('Усть-Каменогорск',  34,  13_328_235,  20_042_050,  1_800_000,  18_242_050,  4_100_000,   536_531, 0.0403),
        ('Павлодар',          27,  13_975_185,  19_465_600,          0,  19_465_600,          0,   720_948, 0.0516),
        ('ИТОГО',            408,  14_071_274, 270_267_339, 12_264_200, 258_003_139, 66_880_100,   632_361, 0.0449),
    ]

    col_x = [0.15, 3.6, 5.1, 6.4, 7.7, 8.9, 10.3, 11.5, 12.6]
    col_w = [3.45, 1.5, 1.3, 1.3, 1.2, 1.4, 1.2, 1.1, 0.7]
    row_h = 0.61

    add_rect(slide, 0.15, 1.25, 13.0, row_h, fill_color=BRAND_COLORS['KIA'])
    hdr_labels = ['Дилерский центр', 'Кол', 'РРЦ', 'Маржа', 'Скидка', 'М. после скидок', 'ДОП', 'Ср. маржа', '%']
    for i, h in enumerate(hdr_labels):
        add_text(slide, h, col_x[i], 1.3, col_w[i], row_h - 0.1,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(dealers):
        y = 1.25 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.15, y, 13.0, row_h - 0.04, fill_color=bg)
        vals = [
            row[0],
            str(row[1]),
            fmt_mln(row[2]),
            fmt_mln(row[3]),
            fmt_mln(row[4]),
            fmt_mln(row[5]),
            fmt_mln(row[6]),
            fmt_num(row[7]),
            fmt_pct(row[8]),
        ]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, col_x[ci], y + 0.07, col_w[ci], row_h - 0.14,
                     font_size=10, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

slide_kia_dealers()

# ═══════════════════════════════════════════════════════════════
# Слайд 4 — KIA анализ по моделям
# ═══════════════════════════════════════════════════════════════
def slide_kia_models():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BRAND_COLORS['KIA'])
    add_text(slide, 'KIA — АНАЛИЗ ПО МОДЕЛЯМ', 0.3, 0.15, 12, 0.8,
             font_size=22, bold=True, color=WHITE)
    add_text(slide, 'Апрель 2026 | 408 авто', 10.0, 0.18, 3.0, 0.5,
             font_size=13, bold=False, color=GOLD, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    models = [
        ('Carnival',         3,  24_323_333,  1_234_417, 0.0508, 0.74),
        ('Ceed',             6,   9_890_000,    437_642, 0.0443, 1.47),
        ('Ceed SW',         10,  11_390_000,    512_550, 0.0450, 2.45),
        ('Cerato',          48,  10_437_917,    382_550, 0.0367, 11.76),
        ('K5',              91,  13_472_418,    675_655, 0.0502, 22.30),
        ('K8',               1,  26_490_000,  1_721_850, 0.0650, 0.25),
        ('Seltos',          26,  12_463_077,    558_340, 0.0448, 6.37),
        ('Soluto',          38,   7_384_737,    298_947, 0.0405, 9.31),
        ('Sorento 2025',    61,  18_811_311,    743_283, 0.0395, 14.95),
        ('Sorento 2026',    11,  22_390_000,  1_123_314, 0.0502, 2.70),
        ('Soul',             4,  12_015_000,    540_675, 0.0450, 0.98),
        ('Sportage',        93,  16_023_976,    788_217, 0.0492, 22.79),
        ('Sportage comfort',16,  12_790_000,    371_200, 0.0290, 3.92),
        ('ИТОГО',          408,  14_071_274,    632_361, 0.0449, 100.0),
    ]

    col_x = [0.2, 3.3, 5.0, 6.8, 8.8, 10.5, 11.8]
    col_w = [3.1, 1.7, 1.8, 2.0, 1.7, 1.3, 1.5]
    row_h = 0.43

    add_rect(slide, 0.2, 1.25, 12.9, row_h, fill_color=BRAND_COLORS['KIA'])
    hdr_labels = ['Модель', 'Кол', 'РРЦ', 'Ср. маржа', 'Доля %', '% авто']
    col_x2 = [0.2, 3.3, 5.0, 6.8, 8.8, 10.5]
    col_w2 = [3.1, 1.7, 1.8, 2.0, 1.7, 2.5]
    for i, h in enumerate(hdr_labels):
        add_text(slide, h, col_x2[i], 1.28, col_w2[i], row_h - 0.05,
                 font_size=11, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(models):
        y = 1.25 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.2, y, 12.9, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_mln(row[2]), fmt_num(row[3]), fmt_pct(row[4]), f"{row[5]:.1f}%"]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, col_x2[ci], y + 0.04, col_w2[ci], row_h - 0.08,
                     font_size=10, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

slide_kia_models()

# ═══════════════════════════════════════════════════════════════
# Слайд 5 — Chevrolet
# ═══════════════════════════════════════════════════════════════
def slide_chevy():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BRAND_COLORS['Chevrolet'])
    add_text(slide, 'CHEVROLET — АНАЛИЗ ПО ДИЛЕРАМ И МОДЕЛЯМ', 0.3, 0.15, 12, 0.8,
             font_size=22, bold=True, color=WHITE)
    add_text(slide, 'Апрель 2026 | 1 787 авто', 9.5, 0.18, 3.5, 0.5,
             font_size=13, bold=False, color=GOLD, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    # Дилеры (левая часть)
    add_rect(slide, 0.15, 1.25, 7.5, 0.45, fill_color=BRAND_COLORS['Chevrolet'])
    add_text(slide, 'Дилерские центры', 0.2, 1.27, 7.4, 0.4,
             font_size=12, bold=True, color=WHITE)

    dealers = [
        ('Актау',          91, 20_700_000, 227_473, 0.0373),
        ('Атырау',         72, 17_100_000, 237_500, 0.0387),
        ('Алматы 1159',     4,  1_200_000, 300_000, 0.0443),
        ('Шымкент 119',   375, 82_340_000, 219_573, 0.0345),
        ('Алматы Zhetysu', 72, 17_450_000, 242_361, 0.0333),
        ('Тараз',         118, 25_950_000, 219_915, 0.0357),
        ('Астана',        108, 29_000_000, 268_519, 0.0360),
        ('Усть-Каменогорск',130,29_500_000,226_923, 0.0369),
        ('Павлодар',       94, 23_100_000, 245_745, 0.0367),
        ('Уральск',       124, 29_050_000, 234_274, 0.0375),
        ('Кызылорда',     130, 28_700_000, 220_769, 0.0359),
        ('Шымкент 935',   469,108_040_000, 230_362, 0.0355),
        ('ИТОГО',        1787,412_130_000, 230_627, 0.0358),
    ]

    dcol_x = [0.2, 3.1, 4.4, 5.8, 7.0]
    dcol_w = [2.9, 1.3, 1.4, 1.2, 0.7]
    row_h = 0.41

    add_rect(slide, 0.15, 1.7, 7.6, row_h, fill_color=RGBColor(0xAA, 0x15, 0x22))
    for i, h in enumerate(['Дилер', 'Кол', 'Маржа', 'Ср.маржа', '%']):
        add_text(slide, h, dcol_x[i], 1.73, dcol_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(dealers):
        y = 1.7 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.15, y, 7.6, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_mln(row[2]), fmt_num(row[3]), fmt_pct(row[4])]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, dcol_x[ci], y + 0.04, dcol_w[ci], row_h - 0.08,
                     font_size=9, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Модели (правая часть)
    add_rect(slide, 8.0, 1.25, 5.15, 0.45, fill_color=BRAND_COLORS['Chevrolet'])
    add_text(slide, 'По моделям', 8.05, 1.27, 5.0, 0.4,
             font_size=12, bold=True, color=WHITE)

    mods = [
        ('Cobalt',      1748, 223_587, 0.0364, 97.8),
        ('Labo',           3, 200_000, 0.0358,  0.2),
        ('Onix',          23, 330_435, 0.0432,  1.3),
        ('Tahoe MY25',    11, 1_118_182, 0.0220, 0.6),
        ('Tracker MY22',   2, 400_000, 0.0381,  0.1),
        ('ИТОГО',       1787, 230_627, 0.0358, 100.0),
    ]
    mcol_x = [8.05, 9.85, 11.0, 12.0, 12.8]
    mcol_w = [1.8,  1.15, 1.0,  0.8,  0.65]
    add_rect(slide, 8.0, 1.7, 5.15, row_h, fill_color=RGBColor(0xAA, 0x15, 0x22))
    for i, h in enumerate(['Модель', 'Кол', 'Ср.маржа', '%', '% авт']):
        add_text(slide, h, mcol_x[i], 1.73, mcol_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(mods):
        y = 1.7 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 8.0, y, 5.15, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_num(row[2]), fmt_pct(row[3]), f"{row[4]:.1f}%"]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, mcol_x[ci], y + 0.04, mcol_w[ci], row_h - 0.08,
                     font_size=10, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

slide_chevy()

# ═══════════════════════════════════════════════════════════════
# Слайд 6 — JAC
# ═══════════════════════════════════════════════════════════════
def slide_jac():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BRAND_COLORS['JAC'])
    add_text(slide, 'JAC — АНАЛИЗ ПО ДИЛЕРАМ И МОДЕЛЯМ', 0.3, 0.15, 12, 0.8,
             font_size=22, bold=True, color=WHITE)
    add_text(slide, 'Апрель 2026 | 88 авто', 10.5, 0.18, 2.5, 0.5,
             font_size=13, bold=False, color=GOLD, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    dealers = [
        ('Шымкент Baidybek',      14,  9_797_143,  2_280_000,  6_535_800,  466_843, 0.0477),
        ('Тараз',                  2, 11_890_000,          0,    951_200,  475_600, 0.0400),
        ('Актау',                  3, 10_923_333,          0,  1_560_600,  520_200, 0.0476),
        ('Атырау',                 5, 10_370_000,     30_000,  2_232_200,  446_440, 0.0431),
        ('Шымкент ONTUSTIK',      11,  8_408_182,    500_000,  3_366_000,  306_000, 0.0364),
        ('Усть-Каменогорск',      12,  9_248_333,  2_350_000,  5_084_800,  423_733, 0.0458),
        ('Шымкент Темирлановское',  7,  7_590_000,  1_250_000,  2_125_200,  303_600, 0.0400),
        ('Астана',                12,  9_183_333,  2_058_720,  4_598_400,  383_200, 0.0417),
        ('Павлодар',               9,  9_701_111,    535_000,  4_192_000,  465_778, 0.0480),
        ('Кызылорда',              6,  8_373_333,    200_000,  2_009_600,  334_933, 0.0400),
        ('Уральск',                7,  8_847_143,    372_800,  2_793_000,  399_000, 0.0451),
        ('ИТОГО',                 88,  9_225_455,  9_576_520, 35_448_800,  402_827, 0.0437),
    ]

    mods = [
        ('J7',      18, 7_828_889, 311_489, 0.0398, 20.45),
        ('J7 PLUS',  2,10_290_000, 411_600, 0.0400,  2.27),
        ('JS4',      5, 8_750_000, 350_000, 0.0400,  5.68),
        ('JS8',      7,11_890_000, 475_600, 0.0400,  7.95),
        ('S3 Pro',  40, 6_855_000, 274_200, 0.0400, 45.45),
        ('T6',       5,12_170_000, 513_680, 0.0422,  5.68),
        ('T8 Pro',   2,15_740_000, 828_700, 0.0526,  2.27),
        ('T9',       9,17_425_556, 971_756, 0.0558, 10.23),
        ('ИТОГО',   88, 9_225_455, 402_827, 0.0437,100.00),
    ]

    # Левая часть — дилеры
    col_x = [0.15, 3.8, 5.1, 6.2, 7.5, 8.7, 9.6]
    col_w = [3.65, 1.3, 1.1, 1.3, 1.2, 0.9, 0.65]
    row_h = 0.46

    add_rect(slide, 0.15, 1.25, 10.1, row_h, fill_color=BRAND_COLORS['JAC'])
    for i, h in enumerate(['Дилер', 'Кол', 'РРЦ', 'Доп', 'Маржа', 'Ср.маржа', '%']):
        add_text(slide, h, col_x[i], 1.28, col_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(dealers):
        y = 1.25 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.15, y, 10.1, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_mln(row[2]), fmt_mln(row[3]), fmt_mln(row[4]), fmt_num(row[5]), fmt_pct(row[6])]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, col_x[ci], y + 0.04, col_w[ci], row_h - 0.08,
                     font_size=9, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    # Правая часть — модели
    mcol_x = [10.4, 11.5, 12.3, 12.9]
    mcol_w = [1.1, 0.8, 0.6, 0.5]
    add_rect(slide, 10.4, 1.25, 2.75, row_h, fill_color=BRAND_COLORS['JAC'])
    for i, h in enumerate(['Модель', 'Кол', '%', '% авт']):
        add_text(slide, h, mcol_x[i], 1.28, mcol_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(mods):
        y = 1.25 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 10.4, y, 2.75, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_pct(row[4]), f"{row[5]:.1f}%"]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, mcol_x[ci], y + 0.04, mcol_w[ci], row_h - 0.08,
                     font_size=9, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

slide_jac()

# ═══════════════════════════════════════════════════════════════
# Слайд 7 — JETOUR
# ═══════════════════════════════════════════════════════════════
def slide_jetour():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=BRAND_COLORS['JETOUR'])
    add_text(slide, 'JETOUR — АНАЛИЗ ПО ДИЛЕРАМ И МОДЕЛЯМ', 0.3, 0.15, 12, 0.8,
             font_size=22, bold=True, color=WHITE)
    add_text(slide, 'Апрель 2026 | 343 авто', 9.8, 0.18, 3.2, 0.5,
             font_size=13, bold=False, color=GOLD, align=PP_ALIGN.RIGHT)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    dealers = [
        ('Алматы Alatau',  26, 11_274_615,  1_250_000, 12_892_100,  495_850, 0.0440),
        ('Актау',          32, 12_208_750,  2_500_000, 20_288_800,  634_025, 0.0519),
        ('Атырау',         34,  8_975_294,  5_395_000, 12_555_700,  369_285, 0.0411),
        ('Алматы Zhetysu', 31, 11_424_194,          0, 17_123_500,  552_371, 0.0484),
        ('Усть-Каменогорск',50,10_170_000,  2_750_000, 25_685_700,  513_714, 0.0505),
        ('Шымкент Ontustik',33,11_050_606,  2_250_000, 16_470_200,  499_097, 0.0452),
        ('Павлодар',        29, 9_593_448,    250_000, 12_856_800,  443_338, 0.0462),
        ('Кызылорда',       22, 9_535_455,  1_400_000,  9_420_100,  428_186, 0.0449),
        ('Тараз',            8, 9_927_500,  1_250_000,  3_776_200,  472_025, 0.0475),
        ('Шымкент Baidybek',29,11_765_862,  2_010_000, 16_666_200,  574_697, 0.0488),
        ('Уральск',         49, 9_592_041,  1_000_000, 22_397_100,  457_084, 0.0477),
        ('ИТОГО',          343,10_480_845, 20_055_000,170_132_400,  496_013, 0.0473),
    ]

    mods = [
        ('Dashing',    9, 10_101_111, 364_044, 0.0360,  2.62),
        ('T1',        30, 15_740_000, 927_733, 0.0589,  8.75),
        ('T2',        29, 16_314_138, 896_103, 0.0549,  8.45),
        ('X50',       84,  7_644_762, 299_600, 0.0392, 24.49),
        ('X70',        1,  8_990_000, 449_500, 0.0500,  0.29),
        ('X70FL',    125,  8_717_680, 417_673, 0.0479, 36.44),
        ('X70Plus',   49, 12_398_163, 548_947, 0.0443, 14.29),
        ('X90Plus',   16, 13_146_250, 519_600, 0.0395,  4.66),
        ('ИТОГО',    343, 10_480_845, 496_013, 0.0473,100.00),
    ]

    col_x = [0.15, 3.8, 5.1, 6.2, 7.5, 8.7, 9.65]
    col_w = [3.65, 1.3, 1.1, 1.3, 1.2, 0.95, 0.65]
    row_h = 0.46

    add_rect(slide, 0.15, 1.25, 10.1, row_h, fill_color=BRAND_COLORS['JETOUR'])
    for i, h in enumerate(['Дилер', 'Кол', 'РРЦ', 'Доп', 'Маржа', 'Ср.маржа', '%']):
        add_text(slide, h, col_x[i], 1.28, col_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(dealers):
        y = 1.25 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.15, y, 10.1, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_mln(row[2]), fmt_mln(row[3]), fmt_mln(row[4]), fmt_num(row[5]), fmt_pct(row[6])]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, col_x[ci], y + 0.04, col_w[ci], row_h - 0.08,
                     font_size=9, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    mcol_x = [10.4, 11.55, 12.2, 12.85]
    mcol_w = [1.15, 0.65, 0.65, 0.5]
    add_rect(slide, 10.4, 1.25, 2.75, row_h, fill_color=BRAND_COLORS['JETOUR'])
    for i, h in enumerate(['Модель', 'Кол', '%', '% авт']):
        add_text(slide, h, mcol_x[i], 1.28, mcol_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(mods):
        y = 1.25 + (ri + 1) * row_h
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 10.4, y, 2.75, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_pct(row[4]), f"{row[5]:.1f}%"]
        bold_r = row[0] == 'ИТОГО'
        for ci, v in enumerate(vals):
            add_text(slide, v, mcol_x[ci], y + 0.04, mcol_w[ci], row_h - 0.08,
                     font_size=9, bold=bold_r, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

slide_jetour()

# ═══════════════════════════════════════════════════════════════
# Слайд 8 — Skoda + Динамика по месяцам
# ═══════════════════════════════════════════════════════════════
def slide_skoda_and_dynamics():
    slide = prs.slides.add_slide(BLANK)
    # Шапка разделена: слева Skoda, справа динамика
    add_rect(slide, 0, 0, 13.33, 1.1, fill_color=DARK_BLUE)
    add_text(slide, 'SKODA + ДИНАМИКА ПО МЕСЯЦАМ', 0.3, 0.15, 12, 0.8,
             font_size=22, bold=True, color=WHITE)
    add_rect(slide, 0, 1.1, 13.33, 0.05, fill_color=GOLD)

    # Skoda блок
    add_rect(slide, 0.15, 1.25, 4.5, 0.4, fill_color=BRAND_COLORS['Skoda'])
    add_text(slide, 'SKODA — Апрель 2026', 0.2, 1.27, 4.4, 0.35,
             font_size=12, bold=True, color=WHITE)

    skoda_data = [
        ('Усть-Каменогорск', 1, 24_690_000, 465_430, 0.0189),
        ('Атырау',           2, 24_690_000, 575_000, 0.0233),
        ('ИТОГО',            3, 24_690_000, 538_477, 0.0218),
    ]
    sc_x = [0.2, 2.4, 3.2, 4.0, 4.6]
    sc_w = [2.2, 0.8, 0.8, 0.6, 0.5]
    add_rect(slide, 0.15, 1.65, 4.5, 0.38, fill_color=BRAND_COLORS['Skoda'])
    for i, h in enumerate(['Дилер', 'Кол', 'Маржа', 'Ср.', '%']):
        add_text(slide, h, sc_x[i], 1.67, sc_w[i], 0.33,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    for ri, row in enumerate(skoda_data):
        y = 1.65 + (ri + 1) * 0.38
        bg = RGBColor(0xFF, 0xF5, 0xCC) if row[0] == 'ИТОГО' else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 0.15, y, 4.5, 0.35, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_mln(row[2]), fmt_num(row[3]), fmt_pct(row[4])]
        for ci, v in enumerate(vals):
            add_text(slide, v, sc_x[ci], y + 0.04, sc_w[ci], 0.3,
                     font_size=9, bold=row[0] == 'ИТОГО', color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

    add_text(slide, 'Модель: Kodiaq NG | 3 авто | Ср. маржа: 538 477 ₸ | 2.18%',
             0.2, 3.1, 4.4, 0.4, font_size=10, bold=False,
             color=DARK_BLUE, align=PP_ALIGN.LEFT)

    # Динамика по месяцам (правая часть)
    add_rect(slide, 5.0, 1.25, 8.15, 0.4, fill_color=ACCENT_BLUE)
    add_text(slide, 'ДИНАМИКА ПО МЕСЯЦАМ — ТОТАЛ (кол-во авто / маржа млн ₸)', 5.05, 1.27, 8.0, 0.35,
             font_size=11, bold=True, color=WHITE)

    months_data = [
        ('Апрель',   1947,  904_596_352, 91_545_700, 996_142_052),
        ('Май',      1836,  784_710_580, 69_158_970, 853_869_550),
        ('Июнь',     2006,  840_067_850, 72_318_000, 912_385_850),
        ('Июль',     2391,  869_755_150, 202_056_300,1_071_811_450),
        ('Август',   1772,  693_722_710, 156_756_500, 850_479_210),
        ('Сентябрь', 1835,  764_821_829, 200_700_260, 965_522_089),
        ('Октябрь',  1983,  839_149_068, 232_227_800,1_071_376_868),
    ]

    mc_x = [5.05, 7.4, 8.7, 9.9, 11.3]
    mc_w = [2.35, 1.3, 1.2, 1.4, 1.85]
    row_h = 0.44

    add_rect(slide, 5.0, 1.65, 8.15, row_h, fill_color=ACCENT_BLUE)
    for i, h in enumerate(['Месяц', 'Кол-во', 'Маржа', 'ДО', 'К получению']):
        add_text(slide, h, mc_x[i], 1.68, mc_w[i], row_h - 0.05,
                 font_size=10, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

    apr_idx = 0  # апрель выделим
    for ri, row in enumerate(months_data):
        y = 1.65 + (ri + 1) * row_h
        is_apr = row[0] == 'Апрель'
        bg = RGBColor(0xFF, 0xF5, 0xCC) if is_apr else (LIGHT_GRAY if ri % 2 == 0 else WHITE)
        add_rect(slide, 5.0, y, 8.15, row_h - 0.03, fill_color=bg)
        vals = [row[0], str(row[1]), fmt_mln(row[2]), fmt_mln(row[3]), fmt_mln(row[4])]
        for ci, v in enumerate(vals):
            add_text(slide, v, mc_x[ci], y + 0.04, mc_w[ci], row_h - 0.08,
                     font_size=10, bold=is_apr, color=DARK_BLUE,
                     align=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)

slide_skoda_and_dynamics()

# ═══════════════════════════════════════════════════════════════
# Слайд 9 — Итоговый
# ═══════════════════════════════════════════════════════════════
def slide_final():
    slide = prs.slides.add_slide(BLANK)
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=DARK_BLUE)
    add_rect(slide, 0, 5.8, 13.33, 0.12, fill_color=GOLD)
    add_rect(slide, 0.5, 1.5, 0.07, 4.0, fill_color=ACCENT_BLUE)

    add_text(slide, 'КЛЮЧЕВЫЕ ВЫВОДЫ', 0.8, 1.4, 11, 0.7,
             font_size=14, bold=False, color=GOLD, align=PP_ALIGN.LEFT)
    add_text(slide, 'АПРЕЛЬ 2026', 0.8, 2.0, 11, 1.0,
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    bullets = [
        '🚗  1 947 авто продано | Маржа к получению: 996 млн ₸',
        '📊  KIA — лидер по средней марже: 605 000 ₸/авто | 364 млн ₸',
        '⭐  JETOUR — высший % маржи: 4.73% | 291 млн ₸',
        '📈  Chevrolet — наибольший объём продаж: 1 787 авто',
        '💰  JAC — 88 авто | 112 млн ₸ к получению',
        '🏆  ДО (допоборудование): 91.5 млн ₸ дополнительный доход',
    ]

    for i, b in enumerate(bullets):
        add_text(slide, b, 0.8, 3.2 + i * 0.55, 11.5, 0.52,
                 font_size=15, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

slide_final()

# ── Сохранить ─────────────────────────────────────────────────
out_path = r"c:\Users\D.Muldabaev\Desktop\Приложения для сервиса\Преза\Маржа_Апрель_2026.pptx"
prs.save(out_path)
print(f"Презентация сохранена: {out_path}")
